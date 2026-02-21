"""
PowerPoint report generator for MetricDiagnostics.

Produces a professional, multi-slide .pptx deck containing:
    1. Title slide
    2. Executive summary (key numbers + model info)
    3. Overall waterfall chart (native PPTX chart)
    4. Variable-level driver summary (bar chart + table)
    5. Per-variable detail slides (distribution chart + rate chart + table)
    6. Feature importance slide
"""

from io import BytesIO
from typing import Dict, List, Optional

import numpy as np
import pandas as pd

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

# ---------------------------------------------------------------------------
# Colour palette
# ---------------------------------------------------------------------------

_COLORS = {
    "dark_bg":   RGBColor(0x1B, 0x1B, 0x2F),   # deep navy
    "accent1":   RGBColor(0x4C, 0x72, 0xB0),   # steel blue  (mix)
    "accent2":   RGBColor(0xDD, 0x84, 0x52),   # burnt orange (rate)
    "base":      RGBColor(0x55, 0xA8, 0x68),   # sage green
    "compare":   RGBColor(0xC4, 0x4E, 0x52),   # muted red
    "predicted": RGBColor(0x81, 0x72, 0xB3),   # purple
    "white":     RGBColor(0xFF, 0xFF, 0xFF),
    "light_grey":RGBColor(0xD0, 0xD0, 0xD0),
    "mid_grey":  RGBColor(0x90, 0x90, 0x90),
    "title_blue":RGBColor(0x1F, 0x3A, 0x5F),   # dark blue for titles
    "header_bg": RGBColor(0x2C, 0x3E, 0x6B),   # table header fill
    "row_alt":   RGBColor(0xEB, 0xEF, 0xF5),   # alternating row
    "positive":  RGBColor(0x27, 0xAE, 0x60),   # green for positive
    "negative":  RGBColor(0xC0, 0x39, 0x2B),   # red for negative
}

# Slide dimensions (standard 13.333 x 7.5 in widescreen)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ---------------------------------------------------------------------------
# Helper functions
# ---------------------------------------------------------------------------

def _set_slide_bg(slide, color: RGBColor = _COLORS["white"]):
    """Set the solid background colour of a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_textbox(slide, left, top, width, height, text,
                 font_size=12, bold=False, color=_COLORS["dark_bg"],
                 alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    """Add a simple text box to a slide."""
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txbox


def _add_table(slide, df: pd.DataFrame, left, top, width, height,
               font_size=9, header_color=_COLORS["header_bg"]):
    """Add a formatted table from a DataFrame to a slide."""
    rows, cols = df.shape
    table_shape = slide.shapes.add_table(rows + 1, cols, left, top, width, height)
    table = table_shape.table

    # Distribute column widths evenly
    col_w = int(width / cols)
    for i in range(cols):
        table.columns[i].width = col_w

    # Header row
    for j, col_name in enumerate(df.columns):
        cell = table.cell(0, j)
        cell.text = str(col_name)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(font_size)
            p.font.bold = True
            p.font.color.rgb = _COLORS["white"]
            p.font.name = "Calibri"
            p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Data rows
    for i in range(rows):
        for j in range(cols):
            cell = table.cell(i + 1, j)
            val = df.iloc[i, j]
            if isinstance(val, float):
                cell.text = f"{val:,.6f}" if abs(val) < 1 else f"{val:,.2f}"
            else:
                cell.text = str(val)

            # Alternating row colour
            if i % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = _COLORS["row_alt"]

            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size)
                p.font.name = "Calibri"
                p.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    return table_shape


def _fmt_pct(val):
    """Format a fraction as a percentage string."""
    if pd.isna(val):
        return "N/A"
    return f"{val:+.1f}%"


def _fmt_val(val, decimals=4):
    """Format a numeric value."""
    if pd.isna(val):
        return "N/A"
    if abs(val) < 1:
        return f"{val:+.{decimals}f}"
    return f"{val:+,.2f}"


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def _slide_title(prs: Presentation, target_col: str, model_type: str):
    """Slide 1 -- Title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _set_slide_bg(slide, _COLORS["dark_bg"])

    _add_textbox(slide,
                 Inches(1), Inches(1.8), Inches(11), Inches(1.5),
                 "Metric Diagnostics Report",
                 font_size=40, bold=True, color=_COLORS["white"],
                 alignment=PP_ALIGN.CENTER)

    _add_textbox(slide,
                 Inches(1), Inches(3.5), Inches(11), Inches(1),
                 f"Target: {target_col}   |   Model: {model_type.upper()}",
                 font_size=20, bold=False, color=_COLORS["light_grey"],
                 alignment=PP_ALIGN.CENTER)

    # Decorative line
    from pptx.shapes.autoshape import Shape
    line = slide.shapes.add_shape(
        1,  # rectangle
        Inches(3), Inches(3.2), Inches(7), Pt(3),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = _COLORS["accent1"]
    line.line.fill.background()


def _slide_executive_summary(prs: Presentation, overall: Dict,
                             n_base: int, n_compare: int,
                             target_col: str, model_type: str,
                             n_features: int):
    """Slide 2 -- Executive summary with key metrics in a card layout."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)

    _add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
                 "Executive Summary",
                 font_size=28, bold=True, color=_COLORS["title_blue"])

    # --- Info banner ---
    info_lines = (
        f"Target Column: {target_col}        "
        f"Model: {model_type.upper()}        "
        f"Features: {n_features}        "
        f"Base n={n_base:,}        "
        f"Compare n={n_compare:,}"
    )
    _add_textbox(slide, Inches(0.5), Inches(1.0), Inches(12), Inches(0.4),
                 info_lines, font_size=12, color=_COLORS["mid_grey"])

    # --- Metric cards (coloured rectangles) ---
    cards = [
        ("Base Metric",    f"{overall['base_metric']:.4f}",    _COLORS["base"]),
        ("Compare Metric", f"{overall['compare_metric']:.4f}", _COLORS["compare"]),
        ("Total Diff",     _fmt_val(overall['total_difference']),  _COLORS["title_blue"]),
        ("Mix Effect",     f"{_fmt_val(overall['distribution_effect'])}  ({_fmt_pct(overall['distribution_pct'])})",
         _COLORS["accent1"]),
        ("Rate Effect",    f"{_fmt_val(overall['relationship_effect'])}  ({_fmt_pct(overall['relationship_pct'])})",
         _COLORS["accent2"]),
        ("Model R-sq",     f"{overall['model_r2_base']:.4f}",  _COLORS["predicted"]),
    ]

    card_w = Inches(3.8)
    card_h = Inches(1.6)
    margin_x = Inches(0.5)
    gap_x = Inches(0.4)
    start_y = Inches(1.8)

    for idx, (label, value, color) in enumerate(cards):
        col = idx % 3
        row = idx // 3
        left = margin_x + col * (card_w + gap_x)
        top = start_y + row * (card_h + Inches(0.3))

        # Card background
        shape = slide.shapes.add_shape(1, left, top, card_w, card_h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        shape.shadow.inherit = False

        # Label
        _add_textbox(slide, left + Inches(0.3), top + Inches(0.2),
                     card_w - Inches(0.6), Inches(0.4),
                     label, font_size=13, bold=False,
                     color=RGBColor(0xFF, 0xFF, 0xFF))
        # Value
        _add_textbox(slide, left + Inches(0.3), top + Inches(0.7),
                     card_w - Inches(0.6), Inches(0.7),
                     value, font_size=22, bold=True,
                     color=RGBColor(0xFF, 0xFF, 0xFF))


def _slide_waterfall(prs: Presentation, overall: Dict):
    """Slide 3 -- Waterfall chart: Base -> Mix -> Rate -> Compare."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)

    _add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
                 "Overall Decomposition (Waterfall)",
                 font_size=28, bold=True, color=_COLORS["title_blue"])

    # python-pptx doesn't have a native waterfall chart type, so we
    # simulate it with a stacked bar chart (invisible base + visible bar).
    chart_data = CategoryChartData()
    chart_data.categories = [
        "Base\nMetric",
        "Distribution\nEffect (Mix)",
        "Relationship\nEffect (Rate)",
        "Compare\nMetric",
    ]

    base_val = overall["base_metric"]
    mix_val = overall["distribution_effect"]
    rate_val = overall["relationship_effect"]
    compare_val = overall["compare_metric"]

    # Invisible base series (the part below the visible bar)
    invisible = [0, base_val, base_val + mix_val, 0]
    # Visible bar series
    visible = [base_val, mix_val, rate_val, compare_val]

    # For negative effects the invisible/visible need adjustment
    for i in [1, 2]:
        if visible[i] < 0:
            invisible[i] = invisible[i] + visible[i]
            visible[i] = abs(visible[i])

    chart_data.add_series("Invisible", invisible)
    chart_data.add_series("Visible", visible)

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_STACKED, Inches(1), Inches(1.2),
        Inches(11), Inches(5.8), chart_data,
    )
    chart = chart_frame.chart
    chart.has_legend = False

    # Make the invisible series truly invisible
    series_inv = chart.series[0]
    series_inv.format.fill.background()
    series_inv.format.line.fill.background()

    # Colour the visible series per category
    series_vis = chart.series[1]
    bar_colors = [_COLORS["base"], _COLORS["accent1"],
                  _COLORS["accent2"], _COLORS["compare"]]
    for idx, color in enumerate(bar_colors):
        pt = series_vis.points[idx]
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = color

    # Data labels on visible series
    series_vis.has_data_labels = True
    dlabels = series_vis.data_labels
    dlabels.show_value = True
    dlabels.font.size = Pt(11)
    dlabels.font.bold = True
    dlabels.font.color.rgb = _COLORS["dark_bg"]
    dlabels.number_format = '0.0000'
    dlabels.label_position = XL_LABEL_POSITION.OUTSIDE_END

    # Style axes
    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.size = Pt(11)
    cat_axis.tick_labels.font.name = "Calibri"
    val_axis = chart.value_axis
    val_axis.tick_labels.font.size = Pt(10)
    val_axis.tick_labels.font.name = "Calibri"
    val_axis.has_title = False

    # Annotation text
    _add_textbox(slide, Inches(1), Inches(7.0), Inches(11), Inches(0.4),
                 f"Total Difference = {_fmt_val(overall['total_difference'])}   |   "
                 f"Mix = {_fmt_val(overall['distribution_effect'])} ({_fmt_pct(overall['distribution_pct'])})   |   "
                 f"Rate = {_fmt_val(overall['relationship_effect'])} ({_fmt_pct(overall['relationship_pct'])})",
                 font_size=11, color=_COLORS["mid_grey"], alignment=PP_ALIGN.CENTER)


def _slide_variable_drivers(prs: Presentation, variable_summary: pd.DataFrame):
    """Slide 4 -- Horizontal bar chart of top variable drivers + table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)

    _add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
                 "Variable-Level Drivers",
                 font_size=28, bold=True, color=_COLORS["title_blue"])

    top = variable_summary.head(10).copy()

    # --- Chart (left half) ---
    chart_data = CategoryChartData()
    chart_data.categories = top["variable"].tolist()
    chart_data.add_series("Mix Effect", top["mix_effect"].tolist())
    chart_data.add_series("Rate Effect", top["rate_effect"].tolist())

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.3), Inches(1.2),
        Inches(6.5), Inches(5.8), chart_data,
    )
    chart = chart_frame.chart
    chart.series[0].format.fill.solid()
    chart.series[0].format.fill.fore_color.rgb = _COLORS["accent1"]
    chart.series[1].format.fill.solid()
    chart.series[1].format.fill.fore_color.rgb = _COLORS["accent2"]
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(10)

    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.size = Pt(9)
    cat_axis.tick_labels.font.name = "Calibri"
    val_axis = chart.value_axis
    val_axis.tick_labels.font.size = Pt(9)

    # --- Table (right half) ---
    display_df = top[["variable", "mix_effect", "rate_effect", "total_effect"]].copy()
    display_df.columns = ["Variable", "Mix Effect", "Rate Effect", "Total Effect"]
    _add_table(slide, display_df,
               Inches(7.2), Inches(1.2), Inches(5.8), Inches(5.0),
               font_size=10)


def _slide_variable_detail(prs: Presentation, var_name: str,
                           detail_df: pd.DataFrame):
    """Slides 5+ -- One slide per variable with distribution chart,
    metric-rate chart, and a detail table."""
    df = detail_df[detail_df["bucket"] != "TOTAL"].copy()
    total_row = detail_df[detail_df["bucket"] == "TOTAL"].iloc[0]
    n_buckets = len(df)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)

    _add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
                 f"Detail: {var_name}",
                 font_size=26, bold=True, color=_COLORS["title_blue"])

    # Subtitle with totals
    _add_textbox(slide, Inches(0.5), Inches(0.85), Inches(12), Inches(0.35),
                 f"Total Mix Effect = {_fmt_val(total_row['mix_effect'])}   |   "
                 f"Total Rate Effect = {_fmt_val(total_row['rate_effect'])}   |   "
                 f"Total Effect = {_fmt_val(total_row['total_effect'])}",
                 font_size=11, color=_COLORS["mid_grey"])

    # --- Left chart: Distribution (base_weight vs compare_weight) ---
    chart_data1 = CategoryChartData()
    chart_data1.categories = df["bucket"].tolist()
    chart_data1.add_series("Base Weight", df["base_weight"].tolist())
    chart_data1.add_series("Compare Weight", df["compare_weight"].tolist())

    cf1 = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.3), Inches(1.3),
        Inches(6.2), Inches(3.0), chart_data1,
    )
    ch1 = cf1.chart
    ch1.series[0].format.fill.solid()
    ch1.series[0].format.fill.fore_color.rgb = _COLORS["base"]
    ch1.series[1].format.fill.solid()
    ch1.series[1].format.fill.fore_color.rgb = _COLORS["compare"]
    ch1.has_legend = True
    ch1.legend.position = XL_LEGEND_POSITION.BOTTOM
    ch1.legend.include_in_layout = False
    ch1.legend.font.size = Pt(8)
    ch1.category_axis.tick_labels.font.size = Pt(7)
    ch1.value_axis.tick_labels.font.size = Pt(8)
    ch1.value_axis.number_format = '0%'

    # Chart title
    ch1.has_title = True
    ch1.chart_title.text_frame.paragraphs[0].text = "Distribution"
    ch1.chart_title.text_frame.paragraphs[0].font.size = Pt(11)
    ch1.chart_title.text_frame.paragraphs[0].font.bold = True

    # --- Right chart: Metric rate (base_rate vs compare_rate vs predicted) ---
    chart_data2 = CategoryChartData()
    chart_data2.categories = df["bucket"].tolist()
    chart_data2.add_series("Base Rate", df["base_rate"].tolist())
    chart_data2.add_series("Compare Rate", df["compare_rate"].tolist())
    chart_data2.add_series("Predicted Compare", df["predicted_compare_rate"].tolist())

    cf2 = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS, Inches(6.8), Inches(1.3),
        Inches(6.2), Inches(3.0), chart_data2,
    )
    ch2 = cf2.chart
    ch2.series[0].format.line.color.rgb = _COLORS["base"]
    ch2.series[0].format.line.width = Pt(2.5)
    ch2.series[0].marker.style = 8    # CIRCLE
    ch2.series[0].marker.size = 8
    ch2.series[1].format.line.color.rgb = _COLORS["compare"]
    ch2.series[1].format.line.width = Pt(2.5)
    ch2.series[1].marker.style = 1    # SQUARE
    ch2.series[1].marker.size = 8
    ch2.series[2].format.line.color.rgb = _COLORS["predicted"]
    ch2.series[2].format.line.width = Pt(2)
    ch2.series[2].format.line.dash_style = 4  # dash
    ch2.series[2].marker.style = 3    # TRIANGLE
    ch2.series[2].marker.size = 7
    ch2.has_legend = True
    ch2.legend.position = XL_LEGEND_POSITION.BOTTOM
    ch2.legend.include_in_layout = False
    ch2.legend.font.size = Pt(8)
    ch2.category_axis.tick_labels.font.size = Pt(7)
    ch2.value_axis.tick_labels.font.size = Pt(8)

    ch2.has_title = True
    ch2.chart_title.text_frame.paragraphs[0].text = "Metric Rate"
    ch2.chart_title.text_frame.paragraphs[0].font.size = Pt(11)
    ch2.chart_title.text_frame.paragraphs[0].font.bold = True

    # --- Bottom table: full detail ---
    table_df = detail_df[[
        "bucket", "base_count", "compare_count",
        "base_weight", "compare_weight", "weight_change",
        "base_rate", "compare_rate", "rate_change",
        "mix_effect", "rate_effect", "total_effect",
    ]].copy()
    table_df.columns = [
        "Bucket", "Base n", "Compare n",
        "Base Wt", "Compare Wt", "Wt Change",
        "Base Rate", "Compare Rate", "Rate Chg",
        "Mix Effect", "Rate Effect", "Total Effect",
    ]
    _add_table(slide, table_df,
               Inches(0.3), Inches(4.5), Inches(12.7), Inches(2.8),
               font_size=8)


def _slide_feature_importance(prs: Presentation, importance_df: pd.DataFrame):
    """Feature importance slide with a horizontal bar chart."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)

    _add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
                 "Feature Importance (from model)",
                 font_size=28, bold=True, color=_COLORS["title_blue"])

    top = importance_df.head(15).copy()

    # Chart
    chart_data = CategoryChartData()
    chart_data.categories = top["feature"].tolist()
    chart_data.add_series("Importance %", top["importance_pct"].tolist())

    cf = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.5), Inches(1.2),
        Inches(7), Inches(5.8), chart_data,
    )
    ch = cf.chart
    ch.series[0].format.fill.solid()
    ch.series[0].format.fill.fore_color.rgb = _COLORS["accent1"]
    ch.has_legend = False

    ch.series[0].has_data_labels = True
    ch.series[0].data_labels.show_value = True
    ch.series[0].data_labels.font.size = Pt(10)
    ch.series[0].data_labels.number_format = '0.0"%"'

    ch.category_axis.tick_labels.font.size = Pt(10)
    ch.value_axis.tick_labels.font.size = Pt(9)

    # Table on the right
    display_df = top[["feature", "importance_pct"]].copy()
    display_df.columns = ["Feature", "Importance %"]
    _add_table(slide, display_df,
               Inches(8), Inches(1.5), Inches(4.5), Inches(5.0),
               font_size=10)


# ---------------------------------------------------------------------------
# Public function
# ---------------------------------------------------------------------------

def generate_report(
    results: Dict,
    target_col: str,
    model_type: str,
    n_base: int,
    n_compare: int,
    n_features: int,
    output_path: str = "metric_diagnostics_report.pptx",
    top_n_variables: int = 5,
) -> str:
    """Generate a PowerPoint report from MetricDiagnostics results.

    Parameters
    ----------
    results : dict
        The dict returned by ``MetricDiagnostics.run()``.
    target_col : str
        Name of the target metric column.
    model_type : str
        Model type used (``'linear'``, ``'logistic'``, ``'lightgbm'``).
    n_base : int
        Number of observations in the base population.
    n_compare : int
        Number of observations in the compare population.
    n_features : int
        Number of feature columns used.
    output_path : str
        File path for output ``.pptx``. Default ``'metric_diagnostics_report.pptx'``.
    top_n_variables : int
        How many top-driver variables to create detail slides for.

    Returns
    -------
    str
        The output file path.
    """
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    overall = results["overall"]
    variable_summary = results["variable_summary"]
    by_variable = results["by_variable"]
    importance_df = results["feature_importance"]

    # 1. Title
    _slide_title(prs, target_col, model_type)

    # 2. Executive summary
    _slide_executive_summary(prs, overall, n_base, n_compare,
                             target_col, model_type, n_features)

    # 3. Waterfall
    _slide_waterfall(prs, overall)

    # 4. Variable drivers
    _slide_variable_drivers(prs, variable_summary)

    # 5. Per-variable detail (top N)
    top_vars = variable_summary.head(top_n_variables)["variable"].tolist()
    for var in top_vars:
        if var in by_variable:
            _slide_variable_detail(prs, var, by_variable[var])

    # 6. Feature importance
    _slide_feature_importance(prs, importance_df)

    prs.save(output_path)
    return output_path
